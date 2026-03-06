#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import logging
import os
import re
import sys
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple

try:
    import requests
    from bs4 import BeautifulSoup
    from dotenv import load_dotenv
    from pptx import Presentation
    from pptx.util import Pt
    from urllib3.exceptions import InsecureRequestWarning
except ModuleNotFoundError as exc:
    missing = str(exc).split("'")[1] if "'" in str(exc) else str(exc)
    print(
        "Ошибка запуска: не хватает python-библиотек.\n"
        f"Не найден модуль: {missing}\n\n"
        "Установите зависимости:\n"
        "python3 -m pip install requests python-dotenv urllib3 beautifulsoup4 python-pptx",
        file=sys.stderr,
    )
    raise SystemExit(2)


load_dotenv()
requests.packages.urllib3.disable_warnings(InsecureRequestWarning)

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


class ConfluenceReleaseStats:
    def __init__(self):
        self.config = {
            "confluence": {
                "url": os.getenv("CONFLUENCE_URL", "https://confluence.sberbank.ru"),
                "token": os.getenv("CONFLUENCE_TOKEN"),
                "space": os.getenv("CONFLUENCE_SPACE", "HRTECH"),
                "source_page_id": os.getenv("SOURCE_PAGE_ID", "18588013525"),
                "verify_ssl": os.getenv("CONFLUENCE_VERIFY_SSL", "false").lower() == "true",
            },
            "jira": {
                "url": os.getenv("JIRA_URL", "https://jira.sberbank.ru"),
                "token": os.getenv("JIRA_TOKEN"),
                "verify_ssl": os.getenv("JIRA_VERIFY_SSL", "false").lower() == "true",
            },
            "pptx": {
                "template_path": os.getenv("PPTX_TEMPLATE_PATH", "/Users/asklimenko/Downloads/ОС ЦРФК 20.02.pptx"),
                "output_path": os.getenv("PPTX_OUTPUT_PATH", ""),
                "slide_index": int(os.getenv("PPTX_SLIDE_INDEX", "1")) - 1,
                "marker_text": os.getenv("PPTX_RELEASES_MARKER", "Релизов не найдено"),
                "section_title": os.getenv("PPTX_SECTION_TITLE", "Мобильные приложения (RN, Android, iOS):"),
            },
        }

        self.allowed_statuses = self.parse_csv_env(
            "RELEASE_STATUSES",
            ["Установлен на ПРОМ", "Готов", "Установка на ПРОМ"],
        )
        self.mobile_keywords = self.parse_csv_env(
            "MOBILE_RELEASE_KEYWORDS",
            ["rn", "react native", "android", "ios", "мп", "mobile"],
        )

        self.headers = {
            "Authorization": f'Bearer {self.config["confluence"]["token"]}',
            "Content-Type": "application/json",
            "Accept": "application/json",
        }
        self.jira_headers = {
            "Authorization": f'Bearer {self.config["jira"]["token"]}',
            "Content-Type": "application/json",
            "Accept": "application/json",
        }

        if os.getenv("DEBUG_MODE", "false").lower() == "true":
            logging.getLogger().setLevel(logging.DEBUG)

    @staticmethod
    def parse_csv_env(name: str, default_values: Sequence[str]) -> List[str]:
        value = os.getenv(name, "")
        if not value.strip():
            return list(default_values)
        return [item.strip() for item in value.split(",") if item.strip()]

    @staticmethod
    def normalize_header_name(value: str) -> str:
        return " ".join(value.strip().lower().replace("\n", " ").split())

    @staticmethod
    def find_column_index(headers: List[str], *aliases: str) -> Optional[int]:
        normalized_aliases = {alias.strip().lower() for alias in aliases}
        for idx, header in enumerate(headers):
            if header in normalized_aliases:
                return idx
        return None

    def get_jira_release_details(self, release_key: str) -> Dict:
        if not self.config["jira"]["token"]:
            logger.info("ℹ️ JIRA_TOKEN не задан, пропускаем обогащение данных для %s", release_key)
            return {}

        try:
            logger.info("🔍 Запрашиваем детали релиза %s из JIRA...", release_key)
            url = f'{self.config["jira"]["url"]}/rest/api/2/issue/{release_key}'
            params = {"fields": "summary,description"}

            response = requests.get(
                url,
                headers=self.jira_headers,
                params=params,
                verify=self.config["jira"]["verify_ssl"],
                timeout=15,
            )

            if response.status_code != 200:
                logger.warning("⚠️ Не удалось получить детали %s: %s", release_key, response.status_code)
                return {}

            data = response.json()
            summary = data["fields"].get("summary", "")
            description = data["fields"].get("description", "")
            services = self.extract_services_from_jira(summary, description)

            logger.info("✅ Получены детали для %s", release_key)
            logger.debug("   Summary: %s", summary)
            logger.debug("   Сервисы: %s", services)

            return {
                "summary": summary,
                "description": description,
                "services": services,
                "full_info": f"{summary} {services}".strip(),
            }
        except Exception as exc:
            logger.error("❌ Ошибка при запросе деталей %s: %s", release_key, exc)
            return {}

    @staticmethod
    def extract_services_from_jira(summary: str, description: str) -> str:
        full_text = f"{summary} {description}" if description else summary
        patterns = [
            r"HRP\.[\w\.-]+\s*\(\d+\)",
            r"CoreUI\s*\(\d+\)",
            r"HRP\.[\w\.-]+\(\d+\)",
            r"[\w\.-]+\s*\(\d{7}\)",
            r"[\w\.-]+\(\d{7}\)",
        ]

        services = []
        for pattern in patterns:
            services.extend(re.findall(pattern, full_text, re.IGNORECASE))

        unique_services = []
        seen = set()
        for service in services:
            if service not in seen:
                unique_services.append(service)
                seen.add(service)

        return ", ".join(unique_services)

    def test_api_connection(self) -> bool:
        url = f'{self.config["confluence"]["url"]}/rest/api/content'
        params = {"limit": 1}

        try:
            logger.info("🔍 Тестируем подключение к Confluence API...")
            response = requests.get(
                url,
                headers=self.headers,
                params=params,
                verify=self.config["confluence"]["verify_ssl"],
                timeout=10,
            )

            logger.info("🔍 Тест API - статус ответа: %s", response.status_code)
            logger.info("🔍 Тест API - размер ответа: %s символов", len(response.text))

            if response.status_code == 200:
                logger.info("✅ Подключение к Confluence API успешно")
                return True
            if response.status_code == 401:
                logger.error("❌ Ошибка авторизации - проверьте токен")
                return False
            if response.status_code == 403:
                logger.error("❌ Недостаточно прав доступа")
                return False

            logger.error("❌ API недоступно: %s", response.status_code)
            logger.error("Ответ сервера: %s", response.text[:500])
            return False
        except requests.exceptions.Timeout:
            logger.error("❌ Превышено время ожидания подключения к API")
            return False
        except requests.exceptions.ConnectionError:
            logger.error("❌ Ошибка подключения к серверу Confluence")
            return False
        except Exception as exc:
            logger.error("❌ Неожиданная ошибка тестирования API: %s", exc)
            return False

    @staticmethod
    def parse_date_from_text(text: str) -> Optional[datetime]:
        for pattern in (r"(\d{4}-\d{2}-\d{2})", r"(\d{2}\.\d{2}\.\d{4})"):
            match = re.search(pattern, text)
            if not match:
                continue
            candidate = match.group(1)
            try:
                if "-" in candidate:
                    return datetime.strptime(candidate, "%Y-%m-%d")
                return datetime.strptime(candidate, "%d.%m.%Y")
            except ValueError:
                continue
        return None

    @staticmethod
    def current_week_range() -> Tuple[datetime, datetime]:
        today = datetime.now()
        days_since_friday = (today.weekday() - 4) % 7
        week_end = today - timedelta(days=days_since_friday)
        week_start = week_end - timedelta(days=6)
        return (
            week_start.replace(hour=0, minute=0, second=0, microsecond=0),
            week_end.replace(hour=23, minute=59, second=59, microsecond=0),
        )

    def get_confluence_page_content(self, page_id: str) -> Tuple[Optional[str], Optional[int]]:
        url = f'{self.config["confluence"]["url"]}/rest/api/content/{page_id}'
        params = {"expand": "body.storage,version"}

        try:
            logger.info("📡 Запрашиваем данные страницы %s...", page_id)
            response = requests.get(
                url,
                headers=self.headers,
                params=params,
                verify=self.config["confluence"]["verify_ssl"],
                timeout=30,
            )

            logger.info("📡 Получен ответ со статусом: %s", response.status_code)

            if response.status_code != 200:
                logger.error("❌ Ошибка HTTP %s при получении страницы %s", response.status_code, page_id)
                logger.error("Ответ сервера: %s", response.text[:1000])
                return None, None

            data = response.json()
            content = data["body"]["storage"]["value"]
            version = data["version"]["number"]

            if not content or not content.strip():
                logger.error("❌ Пустое содержимое страницы %s", page_id)
                return None, None

            logger.info("✅ Получены актуальные данные со страницы %s", page_id)
            logger.info("📄 Размер контента: %s символов", len(content))
            logger.info("🔢 Версия страницы: %s", version)
            logger.info("📝 Первые 200 символов контента: %s", content[:200])
            return content, version
        except requests.exceptions.Timeout:
            logger.error("❌ Превышено время ожидания при получении страницы %s", page_id)
            return None, None
        except requests.exceptions.ConnectionError:
            logger.error("❌ Ошибка подключения при получении страницы %s", page_id)
            return None, None
        except Exception as exc:
            logger.exception("❌ Неожиданная ошибка при получении страницы %s: %s", page_id, exc)
            return None, None

    @staticmethod
    def save_debug_html(content: str, filename: str = "debug_content.html") -> None:
        try:
            with open(filename, "w", encoding="utf-8") as file_obj:
                file_obj.write(content)
            logger.info("💾 HTML контент сохранен в файл: %s", filename)
        except Exception as exc:
            logger.error("❌ Ошибка сохранения HTML: %s", exc)

    def parse_release_table(self, html_content: str) -> List[Dict]:
        releases: List[Dict] = []
        logger.info("=== ПАРСИНГ АКТУАЛЬНЫХ РЕЛИЗОВ С ЗАГРУЗКОЙ ДЕТАЛЕЙ ===")

        try:
            soup = BeautifulSoup(html_content, "html.parser")
            tables = soup.find_all("table")
            logger.info("🔍 Найдено таблиц: %s", len(tables))

            for table_index, table in enumerate(tables):
                rows = table.find_all("tr")
                logger.info("📋 Таблица %s: строк = %s", table_index + 1, len(rows))
                if len(rows) < 2:
                    continue

                header_cells = rows[0].find_all(["th", "td"])
                headers = [self.normalize_header_name(cell.get_text(" ", strip=True)) for cell in header_cells]
                logger.debug("   Заголовки: %s", headers)

                type_idx = self.find_column_index(headers, "тип")
                id_idx = self.find_column_index(headers, "id релиза", "id")
                date_idx = self.find_column_index(headers, "дата", "дата релиза")
                status_idx = self.find_column_index(headers, "статус")
                description_idx = self.find_column_index(headers, "описание релиза", "описание")
                responsible_idx = self.find_column_index(headers, "ответственный")

                if id_idx is None or status_idx is None:
                    logger.debug("   Пропускаем таблицу: не нашли обязательные колонки")
                    continue

                for row_index, row in enumerate(rows[1:], 1):
                    cells = row.find_all(["td", "th"])
                    if len(cells) <= max(id_idx, status_idx):
                        continue

                    try:
                        release_type = (
                            cells[type_idx].get_text(" ", strip=True)
                            if type_idx is not None and len(cells) > type_idx
                            else ""
                        )
                        release_cell = cells[id_idx] if len(cells) > id_idx else None
                        release_date = (
                            cells[date_idx].get_text(" ", strip=True)
                            if date_idx is not None and len(cells) > date_idx
                            else ""
                        )
                        status = cells[status_idx].get_text(" ", strip=True) if len(cells) > status_idx else ""
                        responsible = (
                            cells[responsible_idx].get_text(" ", strip=True)
                            if responsible_idx is not None and len(cells) > responsible_idx
                            else ""
                        )
                        original_description = (
                            cells[description_idx].get_text(" ", strip=True)
                            if description_idx is not None and len(cells) > description_idx
                            else ""
                        )
                        release_full_text = release_cell.get_text(" ", strip=True) if release_cell else ""

                        release_link = ""
                        release_id = ""
                        if release_cell:
                            link_tag = release_cell.find("a", href=True)
                            if link_tag:
                                release_link = link_tag["href"]
                            hrp_match = re.search(r"HRPRELEASE-(\d+)", release_full_text, re.IGNORECASE)
                            if hrp_match:
                                release_id = hrp_match.group(1)

                        if status not in self.allowed_statuses or not release_id:
                            continue

                        release_datetime = self.parse_date_from_text(release_date)
                        jira_release_key = f"HRPRELEASE-{release_id}"
                        jira_details = self.get_jira_release_details(jira_release_key)

                        if jira_details and jira_details.get("services"):
                            extended_description = f'{jira_details.get("summary", "")} {jira_details.get("services", "")}'.strip()
                        elif jira_details and jira_details.get("summary"):
                            extended_description = jira_details["summary"]
                        elif original_description:
                            extended_description = original_description
                        else:
                            desc_without_id = re.sub(r"^HRPRELEASE-\d+\s*", "", release_full_text, flags=re.IGNORECASE)
                            desc_without_id = re.sub(r"^[\s\-–—]+", "", desc_without_id)
                            extended_description = desc_without_id.strip()

                        release_info = {
                            "id": release_id,
                            "key": jira_release_key,
                            "type": release_type,
                            "date": release_date,
                            "datetime": release_datetime,
                            "status": status,
                            "responsible": responsible,
                            "original_description": original_description,
                            "full_description": extended_description,
                            "jira_link": release_link,
                            "full_text": release_full_text,
                            "jira_details": jira_details,
                        }

                        releases.append(release_info)
                        logger.info("✅ %s: %s", jira_release_key, status)
                        logger.debug("   Описание: %s", extended_description)
                    except Exception as exc:
                        logger.error("❌ Ошибка при парсинге строки %s: %s", row_index, exc)
                        continue
        except Exception as exc:
            logger.exception("❌ Критическая ошибка при парсинге HTML: %s", exc)

        logger.info("📊 Всего найдено релизов для анализа: %s", len(releases))
        return releases

    def filter_weekly_mobile_releases(self, releases: List[Dict]) -> List[Dict]:
        week_start, week_end = self.current_week_range()
        logger.info("📅 Фильтруем релизы за неделю: %s - %s", week_start.strftime("%d.%m.%Y"), week_end.strftime("%d.%m.%Y"))
        logger.info("📌 Ключевые слова мобильных релизов: %s", ", ".join(self.mobile_keywords))

        filtered = []
        for release in releases:
            release_date = release.get("datetime")
            if not release_date:
                logger.debug("⏭️ %s пропущен: не определена дата", release["key"])
                continue
            if not (week_start <= release_date <= week_end):
                logger.debug("⏭️ %s пропущен: не входит в текущую неделю", release["key"])
                continue

            search_text = " ".join(
                [
                    release.get("type", ""),
                    release.get("status", ""),
                    release.get("full_description", ""),
                    release.get("original_description", ""),
                    release.get("full_text", ""),
                    release.get("jira_details", {}).get("summary", ""),
                ]
            ).lower()

            if self.mobile_keywords and not any(keyword.lower() in search_text for keyword in self.mobile_keywords):
                logger.debug("⏭️ %s пропущен: не похож на мобильный релиз", release["key"])
                continue

            logger.info("📌 Включаем в презентацию: %s", release["key"])
            filtered.append(release)

        filtered.sort(key=lambda item: item["datetime"] or datetime.min, reverse=True)
        logger.info("📊 Мобильных релизов за неделю найдено: %s", len(filtered))
        return filtered

    def build_slide_lines(self, releases: List[Dict]) -> List[str]:
        lines = [self.config["pptx"]["section_title"], "--"]
        if not releases:
            lines.append("Релизов не найдено")
            return lines

        for release in releases:
            date_label = release["datetime"].strftime("%d.%m") if release.get("datetime") else (release.get("date") or "без даты")
            description = release.get("full_description") or "Описание не найдено"
            responsible = f" ({release['responsible']})" if release.get("responsible") else ""
            lines.append(f"- {release['key']} [{date_label}, {release['status']}] - {description}{responsible}")
        return lines

    def find_target_shape(self, slide):
        marker_text = self.config["pptx"]["marker_text"]
        section_title = self.config["pptx"]["section_title"]

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_value = shape.text or ""
            if marker_text and marker_text in text_value:
                logger.info("🎯 Найден текстовый блок по маркеру '%s'", marker_text)
                return shape
            if section_title and section_title.split(":")[0] in text_value:
                logger.info("🎯 Найден текстовый блок по заголовку секции")
                return shape

        raise RuntimeError(
            "Не удалось найти текстовый блок для релизов. "
            "Проверьте PPTX_RELEASES_MARKER или шаблон презентации."
        )

    def update_presentation(self, releases: List[Dict]) -> Path:
        template_path = Path(self.config["pptx"]["template_path"])
        output_path = Path(
            self.config["pptx"]["output_path"]
            or str(template_path.with_name(f"{template_path.stem} (авто).pptx"))
        )

        if not template_path.exists():
            raise FileNotFoundError(f"Шаблон презентации не найден: {template_path}")

        logger.info("📝 Открываем презентацию: %s", template_path)
        presentation = Presentation(str(template_path))
        slide_index = self.config["pptx"]["slide_index"]

        if slide_index < 0 or slide_index >= len(presentation.slides):
            raise RuntimeError(f"Некорректный номер слайда: {slide_index + 1}")

        slide = presentation.slides[slide_index]
        target_shape = self.find_target_shape(slide)
        lines = self.build_slide_lines(releases)
        text_frame = target_shape.text_frame

        logger.info("🧩 Заполняем первый слайд: %s строк", len(lines))
        text_frame.clear()
        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = Pt(14 if index < 2 else 12)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))
        logger.info("✅ Презентация сохранена: %s", output_path)
        return output_path

    def generate_weekly_presentation(self) -> bool:
        logger.info("🚀 Начинаем генерацию weekly-отчета для презентации")

        if not self.test_api_connection():
            logger.error("❌ Проблемы с подключением к Confluence API - прерываем выполнение")
            return False

        source_content, _ = self.get_confluence_page_content(self.config["confluence"]["source_page_id"])
        if not source_content:
            logger.error("❌ Не удалось получить актуальное содержимое исходной страницы")
            return False

        self.save_debug_html(source_content)
        releases = self.parse_release_table(source_content)
        if not releases:
            logger.error("❌ После парсинга не найдено ни одного релиза")
            return False

        weekly_mobile_releases = self.filter_weekly_mobile_releases(releases)
        self.update_presentation(weekly_mobile_releases)

        logger.info("🎉 Weekly-отчет для презентации успешно сформирован")
        return True


def main() -> int:
    logger.info("▶️ Старт confluence_to_pptx.py")

    required_vars = ["CONFLUENCE_TOKEN"]
    missing_vars = [var for var in required_vars if not os.getenv(var)]
    if missing_vars:
        logger.error("❌ Отсутствуют переменные окружения: %s", ", ".join(missing_vars))
        logger.info("Добавьте в .env минимум:")
        logger.info("CONFLUENCE_TOKEN=ваш_токен")
        logger.info("PPTX_TEMPLATE_PATH=/путь/к/шаблону.pptx")
        logger.info("SOURCE_PAGE_ID=18588013525")
        return 1

    try:
        reporter = ConfluenceReleaseStats()
        return 0 if reporter.generate_weekly_presentation() else 1
    except Exception as exc:
        logger.exception("💥 Критическая ошибка: %s", exc)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
